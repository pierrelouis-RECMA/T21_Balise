import io
import copy
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- CONFIGURATION ---
SLIDE_INDEX_DETAILS = 5  # La Slide 06 (index 5) est le modèle pour les détails
MAX_AGENTS_PER_DETAIL_SLIDE = 4 # Capacité de votre design "Glass" par page

def format_nbb(val):
    """Formate les montants : signe, arrondi à 1 décimale et suffixe 'm'."""
    if pd.isna(val) or val == 0:
        return "0m"
    sign = "+" if val > 0 else ""
    return f"{sign}{val:.1f}m"

def duplicate_slide(prs, source_slide):
    """Duplique une slide en conservant exactement le layout et les formes."""
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

def replace_tags_in_shape(shape, replacements):
    """Parcourt les textes d'une forme (y compris les groupes) pour remplacer les balises."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            replace_tags_in_shape(s, replacements)
    elif shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for tag, value in replacements.items():
                    if tag in run.text:
                        run.text = run.text.replace(tag, str(value))

def build_agency_pptx(df_excel, template_path):
    # 1. NETTOYAGE ET CALCULS
    df = df_excel.copy()
    df['Integrated Spends'] = pd.to_numeric(df['Integrated Spends'], errors='coerce').fillna(0)
    
    # Pivot pour le classement NBB
    summary = df.groupby('Agency').agg({'Integrated Spends': 'sum'}).reset_index()
    summary = summary.sort_values('Integrated Spends', ascending=False).reset_index(drop=True)

    # 2. CHARGEMENT DU TEMPLATE
    prs = Presentation(template_path)
    
    # 3. PRÉPARATION DES REMPLACEMENTS (Slides 02 à 05)
    replacements = {}
    for i, row in summary.head(14).iterrows():
        idx = i + 1
        ag_name = row['Agency']
        replacements[f"{{{{AG_{idx}}}}}"] = ag_name.upper()
        replacements[f"{{{{NBB_{idx}}}}}"] = format_nbb(row['Integrated Spends'])
        replacements[f"{{{{RANK_{idx}}}}}"] = str(idx)

        # Concaténation des Top Wins et Departures (>3$m)
        ag_data = df[df['Agency'] == ag_name]
        wins = ag_data[ag_data['NewBiz'] == 'WIN'].sort_values('Integrated Spends', ascending=False).head(3)
        deps = ag_data[ag_data['NewBiz'] == 'DEPARTURE'].sort_values('Integrated Spends', ascending=True).head(3)
        rets = ag_data[ag_data['NewBiz'] == 'RETENTION'].head(3)

        replacements[f"{{{{TOPWINS_{idx}}}}}"] = " · ".join([f"{r.Advertiser} {format_nbb(r['Integrated Spends'])}" for r in wins.itertuples()])
        replacements[f"{{{{TOPDEPS_{idx}}}}}"] = " · ".join([f"{r.Advertiser} {format_nbb(r['Integrated Spends'])}" for r in deps.itertuples()])
        replacements[f"{{{{TOPRET_{idx}}}}}"] = " · ".