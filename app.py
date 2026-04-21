"""
app.py — NBB Report Generator
──────────────────────────────
Upload an Excel → get a filled PPTX.

Local:
    pip install flask pandas openpyxl python-pptx lxml
    python app.py

Render.com:
    Start command : gunicorn app:app
    Build command : pip install -r requirements.txt
"""

import io, json, math, os, threading, traceback, uuid
from datetime import datetime, timedelta

import pandas as pd
from flask import Flask, abort, jsonify, render_template_string, request, send_file
from pptx import Presentation

from fill_template import load_data_from_df, build_placeholders, replace_all_placeholders
from generate_pptx_v2 import build_agency_pptx

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024  # 20 MB

# ── Template path (same directory as app.py) ──────────────────
TEMPLATE_NAME = "T21_HK_Agencies_Glass_v13.pptx"
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), TEMPLATE_NAME)

# ── In-memory file cache (10-min TTL) ─────────────────────────
_cache: dict = {}
_lock = threading.Lock()

def store_file(data: bytes, filename: str) -> str:
    token  = str(uuid.uuid4())
    expiry = datetime.now() + timedelta(minutes=10)
    with _lock:
        _cache[token] = (data, filename, expiry)
    return token

def purge_cache():
    now = datetime.now()
    with _lock:
        dead = [k for k,(_, _, exp) in _cache.items() if exp < now]
        for k in dead:
            del _cache[k]

# ─────────────────────────────────────────────────────────────
# HTML UI
# ─────────────────────────────────────────────────────────────
HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>NBB Report Generator</title>
<link href="https://fonts.googleapis.com/css2?family=DM+Mono:wght@400;500&family=DM+Sans:wght@300;400;500;600&display=swap" rel="stylesheet">
<style>
:root {
  --bg:#0A0E1A; --surface:#111827; --surface2:#1C2333; --border:#1E293B;
  --accent:#38BDF8; --win:#10B981; --dep:#F43F5E;
  --text:#E2E8F0; --muted:#64748B;
  --mono:'DM Mono',monospace; --sans:'DM Sans',sans-serif;
}
*,*::before,*::after{box-sizing:border-box;margin:0;padding:0}
body{background:var(--bg);color:var(--text);font-family:var(--sans);min-height:100vh;display:flex;flex-direction:column;align-items:center}
header{width:100%;padding:18px 40px;border-bottom:1px solid var(--border);display:flex;align-items:center;gap:12px;background:var(--surface)}
.dot{width:10px;height:10px;border-radius:50%;background:var(--accent);box-shadow:0 0 12px var(--accent)}
header h1{font-size:14px;font-weight:500;letter-spacing:.12em;text-transform:uppercase}
header span{margin-left:auto;font-family:var(--mono);font-size:11px;color:var(--muted)}
main{width:100%;max-width:760px;padding:56px 24px 80px;display:flex;flex-direction:column;gap:28px}
.tagline{text-align:center}
.tagline h2{font-size:30px;font-weight:300;letter-spacing:-.02em;color:#fff}
.tagline h2 em{font-style:normal;color:var(--accent)}
.tagline p{margin-top:8px;font-size:13px;color:var(--muted);line-height:1.6}
.card{background:var(--surface);border:1px solid var(--border);border-radius:12px;padding:26px 30px}
.card-title{font-size:11px;font-weight:500;letter-spacing:.14em;text-transform:uppercase;color:var(--muted);margin-bottom:18px;display:flex;align-items:center;gap:8px}
.card-title::before{content:'';display:block;width:18px;height:1px;background:var(--accent)}
.form-row{display:grid;grid-template-columns:1fr 1fr;gap:14px;margin-bottom:14px}
.form-group{display:flex;flex-direction:column;gap:5px}
.form-group label{font-size:11px;color:var(--muted);font-weight:500;letter-spacing:.05em}
.form-group select,.form-group input[type=text]{background:var(--surface2);border:1px solid var(--border);border-radius:7px;padding:10px 13px;color:var(--text);font-family:var(--sans);font-size:14px;outline:none;transition:border-color .2s;-webkit-appearance:none;appearance:none}
.form-group select:focus,.form-group input:focus{border-color:var(--accent)}
.dropzone{border:1.5px dashed var(--border);border-radius:10px;padding:32px 20px;text-align:center;cursor:pointer;transition:border-color .2s,background .2s;position:relative;overflow:hidden}
.dropzone:hover,.dropzone.drag-over{border-color:var(--accent);background:rgba(56,189,248,.04)}
.dropzone input[type=file]{position:absolute;inset:0;opacity:0;cursor:pointer;width:100%;height:100%}
.dz-icon{font-size:26px;margin-bottom:8px;opacity:.6}
.dz-label{font-size:14px;font-weight:500}
.dz-hint{font-size:11px;color:var(--muted);margin-top:3px;font-family:var(--mono)}
.dz-filename{display:none;margin-top:10px;padding:7px 12px;background:rgba(16,185,129,.1);border:1px solid rgba(16,185,129,.3);border-radius:6px;font-size:12px;font-family:var(--mono);color:var(--win)}
.col-map{display:none;flex-direction:column;gap:0;margin-top:18px}
.col-map.visible{display:flex}
.col-map-title{font-size:11px;letter-spacing:.1em;text-transform:uppercase;color:var(--muted);margin-bottom:10px}
.col-row{display:grid;grid-template-columns:180px 20px 1fr;align-items:center;gap:8px;padding:7px 0;border-bottom:1px solid var(--border)}
.col-row:last-child{border-bottom:none}
.col-row label{font-size:11px;color:var(--muted);font-family:var(--mono)}
.col-row .arrow{color:var(--border);font-size:13px;text-align:center}
.col-row select{background:var(--surface2);border:1px solid var(--border);border-radius:6px;padding:6px 9px;color:var(--text);font-family:var(--mono);font-size:11px;outline:none;transition:border-color .2s;-webkit-appearance:none;appearance:none}
.col-row select:focus{border-color:var(--accent)}
.col-row select.ok{border-color:rgba(16,185,129,.5);color:var(--win)}
.col-row select.err{border-color:rgba(244,63,94,.5);color:var(--dep)}
.btn{width:100%;margin-top:22px;padding:14px;background:var(--accent);color:#0A0E1A;border:none;border-radius:8px;font-family:var(--sans);font-size:15px;font-weight:600;letter-spacing:.02em;cursor:pointer;transition:background .2s,opacity .2s;display:flex;align-items:center;justify-content:center;gap:8px}
.btn:hover:not(:disabled){background:#7DD3FC}
.btn:disabled{opacity:.4;cursor:not-allowed}
.status{display:none;padding:14px 18px;border-radius:9px;font-size:13px;line-height:1.6;margin-top:14px}
.status.loading{display:flex;align-items:center;gap:12px;background:rgba(56,189,248,.07);border:1px solid rgba(56,189,248,.2);color:var(--accent)}
.status.success{display:block;background:rgba(16,185,129,.07);border:1px solid rgba(16,185,129,.25);color:var(--win)}
.status.error{display:block;background:rgba(244,63,94,.07);border:1px solid rgba(244,63,94,.25);color:var(--dep)}
.spinner{width:17px;height:17px;border:2px solid rgba(56,189,248,.2);border-top-color:var(--accent);border-radius:50%;animation:spin .7s linear infinite;flex-shrink:0}
@keyframes spin{to{transform:rotate(360deg)}}
.dl-link{display:inline-block;margin-top:10px;padding:9px 18px;background:var(--win);color:#fff;text-decoration:none;border-radius:6px;font-weight:600;font-size:13px;transition:opacity .2s}
.dl-link:hover{opacity:.85}
.info-grid{display:grid;grid-template-columns:1fr 1fr 1fr;gap:10px}
.info-item{background:var(--surface2);border:1px solid var(--border);border-radius:8px;padding:13px 15px}
.info-label{font-size:10px;letter-spacing:.1em;text-transform:uppercase;color:var(--muted);margin-bottom:3px}
.info-val{font-size:20px;font-weight:600;font-family:var(--mono)}
.info-sub{font-size:11px;color:var(--muted);margin-top:1px}
.req-cols{display:flex;flex-direction:column;gap:5px}
.req-row{display:flex;align-items:baseline;gap:10px;padding:8px 11px;background:var(--surface2);border-radius:6px}
.req-name{font-family:var(--mono);font-size:11px;color:var(--accent);min-width:190px;flex-shrink:0}
.req-desc{color:var(--muted);font-size:11px}
.req-ex{margin-left:auto;font-family:var(--mono);font-size:10px;color:var(--muted);background:var(--border);padding:2px 6px;border-radius:3px;flex-shrink:0}
.badge{font-size:9px;padding:2px 5px;border-radius:3px;font-weight:600;letter-spacing:.05em;text-transform:uppercase}
.badge.r{background:rgba(244,63,94,.15);color:var(--dep)}
.badge.o{background:rgba(100,116,139,.15);color:var(--muted)}
</style>
</head>
<body>
<header>
  <div class="dot"></div>
  <h1>NBB Report Generator</h1>
  <span>Powered by Render · PPTX</span>
</header>
<main>
  <div class="tagline">
    <h2>Upload your Excel,<br>get your <em>presentation</em>.</h2>
    <p>Static slides (Key Findings, TOP moves, Agencies, Groups, Retentions) are auto-filled.<br>
       Agency detail cards are generated dynamically — 4 agencies per slide.</p>
  </div>

  <div class="card">
    <div class="card-title">Required Excel columns</div>
    <div class="req-cols">
      <div class="req-row"><span class="req-name">Agency</span><span class="req-desc">Agency name</span><span class="req-ex">MINDSHARE</span><span class="badge r">required</span></div>
      <div class="req-row"><span class="req-name">NewBiz</span><span class="req-desc">WIN · DEPARTURE · RETENTION</span><span class="req-ex">WIN</span><span class="badge r">required</span></div>
      <div class="req-row"><span class="req-name">Advertiser</span><span class="req-desc">Brand / client name</span><span class="req-ex">RECKITT</span><span class="badge r">required</span></div>
      <div class="req-row"><span class="req-name">Integrated Spends</span><span class="req-desc">Budget $m — positive WIN, negative DEPARTURE</span><span class="req-ex">+150.0</span><span class="badge r">required</span></div>
      <div class="req-row"><span class="req-name">Date of announcement</span><span class="req-desc">Shown on agency card</span><span class="req-ex">2025-01-01</span><span class="badge o">optional</span></div>
    </div>
  </div>

  <div class="card">
    <div class="card-title">Generate report</div>
    <div class="form-row">
      <div class="form-group">
        <label>Market</label>
        <select id="market">
          <option value="Mexico">Mexico</option>
          <option value="Hong_Kong">Hong Kong</option>
          <option value="Singapore">Singapore</option>
          <option value="Indonesia">Indonesia</option>
          <option value="Other">Other</option>
        </select>
      </div>
      <div class="form-group">
        <label>Year</label>
        <input type="text" id="year" value="2025" placeholder="2025">
      </div>
    </div>

    <div class="dropzone" id="dz">
      <input type="file" id="fi" accept=".xlsx,.xls">
      <div class="dz-icon">📊</div>
      <div class="dz-label">Drop your Excel file here</div>
      <div class="dz-hint">.xlsx or .xls · max 20 MB</div>
      <div class="dz-filename" id="dzName"></div>
    </div>

    <div class="col-map" id="colMap">
      <div class="col-map-title">Column mapping</div>
    </div>

    <button class="btn" id="btn" disabled>Generate Presentation</button>
    <div class="status loading" id="sLoad"><div class="spinner"></div><span>Generating…</span></div>
    <div class="status success" id="sOk"></div>
    <div class="status error"   id="sErr"></div>
  </div>

  <div class="card" id="statsCard" style="display:none">
    <div class="card-title">Last generation</div>
    <div class="info-grid">
      <div class="info-item"><div class="info-label">Agencies</div><div class="info-val" id="stAg">—</div><div class="info-sub">in Excel</div></div>
      <div class="info-item"><div class="info-label">Slides</div><div class="info-val" id="stSl">—</div><div class="info-sub">total slides</div></div>
      <div class="info-item"><div class="info-label">Records</div><div class="info-val" id="stRe">—</div><div class="info-sub">W+D+R rows</div></div>
    </div>
  </div>
</main>

<script src="https://cdnjs.cloudflare.com/ajax/libs/xlsx/0.18.5/xlsx.full.min.js"></script>
<script>
const FIELDS = [
  {key:'Agency',required:true},{key:'NewBiz',required:true},
  {key:'Advertiser',required:true},{key:'Integrated Spends',required:true},
  {key:'Date of announcement',required:false},{key:'Incumbent',required:false},
];
let detectedCols = [];
const dz = document.getElementById('dz');
const fi = document.getElementById('fi');

dz.addEventListener('dragover', e=>{e.preventDefault();dz.classList.add('drag-over')});
dz.addEventListener('dragleave', ()=>dz.classList.remove('drag-over'));
dz.addEventListener('drop', e=>{e.preventDefault();dz.classList.remove('drag-over');if(e.dataTransfer.files[0]){fi.files=e.dataTransfer.files;handle(e.dataTransfer.files[0])}});
fi.addEventListener('change', ()=>{if(fi.files[0])handle(fi.files[0])});

function handle(file){
  document.getElementById('dzName').style.display='block';
  document.getElementById('dzName').textContent='✓ '+file.name;
  const r=new FileReader();
  r.onload=e=>{
    try{
      const wb=XLSX.read(new Uint8Array(e.target.result),{type:'array',sheetRows:2});
      detectedCols=(XLSX.utils.sheet_to_json(wb.Sheets[wb.SheetNames[0]],{header:1,defval:''})[0]||[]).map(String);
      buildMap(detectedCols);
      document.getElementById('btn').disabled=false;
    }catch{document.getElementById('btn').disabled=false}
  };
  r.readAsArrayBuffer(file);
}

function buildMap(cols){
  const map=document.getElementById('colMap');
  map.innerHTML='<div class="col-map-title">Column mapping — match to expected fields</div>';
  let needsMap=false;
  FIELDS.forEach(f=>{
    if(!cols.includes(f.key))needsMap=true;
    const row=document.createElement('div');row.className='col-row';
    const lbl=document.createElement('label');lbl.textContent=f.key;
    const arrow=document.createElement('div');arrow.className='arrow';arrow.textContent='→';
    const sel=document.createElement('select');sel.dataset.field=f.key;
    const blank=document.createElement('option');blank.value='';blank.textContent=f.required?'— select —':'— skip —';
    sel.appendChild(blank);
    cols.forEach(c=>{const o=document.createElement('option');o.value=c;o.textContent=c;if(c===f.key)o.selected=true;sel.appendChild(o)});
    sel.className=cols.includes(f.key)?'ok':(f.required?'err':'');
    sel.addEventListener('change',()=>sel.className=sel.value?'ok':(f.required?'err':''));
    row.append(lbl,arrow,sel);map.appendChild(row);
  });
  map.classList.toggle('visible',needsMap||true);
}

function getMap(){
  const m={};
  document.querySelectorAll('#colMap select[data-field]').forEach(s=>{m[s.dataset.field]=s.value||s.dataset.field});
  return m;
}

document.getElementById('btn').addEventListener('click',async()=>{
  if(!fi.files[0])return;
  const [sLoad,sOk,sErr,btn]=['sLoad','sOk','sErr','btn'].map(id=>document.getElementById(id));
  sLoad.style.display='flex';sOk.style.display='none';sErr.style.display='none';btn.disabled=true;
  const fd=new FormData();
  fd.append('file',fi.files[0]);
  fd.append('market',document.getElementById('market').value);
  fd.append('year',document.getElementById('year').value);
  fd.append('col_map',JSON.stringify(getMap()));
  try{
    const res=await fetch('/generate',{method:'POST',body:fd});
    const d=await res.json();
    sLoad.style.display='none';btn.disabled=false;
    if(d.status==='success'){
      sOk.style.display='block';
      sOk.innerHTML=`<strong>✅ Ready!</strong><br>${d.filename}<br><a class="dl-link" href="${d.download_url}" download>⬇ Download PPTX</a>`;
      document.getElementById('statsCard').style.display='block';
      document.getElementById('stAg').textContent=d.agencies;
      document.getElementById('stSl').textContent=d.total_slides;
      document.getElementById('stRe').textContent=d.records;
    }else{
      sErr.style.display='block';
      sErr.innerHTML=`<strong>❌ Error</strong><br>${d.error}`;
    }
  }catch(e){
    sLoad.style.display='none';btn.disabled=false;
    sErr.style.display='block';sErr.innerHTML=`<strong>❌ Network error</strong><br>${e.message}`;
  }
});
</script>
</body>
</html>"""

# ─────────────────────────────────────────────────────────────
# ROUTES
# ─────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template_string(HTML)

@app.route("/generate", methods=["POST"])
def generate():
    try:
        file    = request.files.get("file")
        market  = request.form.get("market", "Unknown")
        year    = request.form.get("year",   "2025")
        col_map = json.loads(request.form.get("col_map", "{}"))

        if not file:
            return jsonify({"status": "error", "error": "No file uploaded"}), 400

        # ── Load & remap columns ──────────────────────────────
        df = pd.read_excel(file)
        rename = {v: k for k, v in col_map.items() if v != k and v in df.columns}
        if rename:
            df = df.rename(columns=rename)

        missing = [c for c in ["Agency","NewBiz","Advertiser","Integrated Spends"]
                   if c not in df.columns]
        if missing:
            return jsonify({
                "status": "error",
                "error": f"Missing columns: {', '.join(missing)}. "
                         f"Available: {', '.join(df.columns.tolist())}"
            }), 400

        if not os.path.exists(TEMPLATE_PATH):
            return jsonify({
                "status": "error",
                "error": f"Template not found: {TEMPLATE_NAME}"
            }), 500

        # ── Step 1: fill static slides 1–6 ───────────────────
        data = load_data_from_df(df)
        ph   = build_placeholders(data)
        prs  = Presentation(TEMPLATE_PATH)
        replace_all_placeholders(prs, ph)
        buf = io.BytesIO()
        prs.save(buf)
        prefilled = buf.getvalue()

        # ── Step 2: generate agency card slides 7+ ────────────
        pptx_bytes = build_agency_pptx(df, TEMPLATE_PATH,
                                       prefilled_prs_bytes=prefilled)

        # ── Compute stats ─────────────────────────────────────
        import zipfile, re as _re
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
            total_slides = len([n for n in z.namelist()
                                if _re.match(r"ppt/slides/slide\d+\.xml$", n)])

        n_agencies = int(df["Agency"].dropna().nunique())
        ts         = datetime.now().strftime("%Y%m%d_%H%M")
        filename   = f"NBB_{market}_{year}_{ts}.pptx"
        token      = store_file(pptx_bytes, filename)
        dl_url     = f"{request.host_url.rstrip('/')}/download/{token}"

        purge_cache()

        return jsonify({
            "status":       "success",
            "download_url": dl_url,
            "filename":     filename,
            "agencies":     n_agencies,
            "total_slides": total_slides,
            "records":      len(df),
            "expires_in":   "10 minutes",
        })

    except Exception as e:
        return jsonify({
            "status": "error",
            "error":  str(e),
            "detail": traceback.format_exc()
        }), 500


@app.route("/download/<token>")
def download(token):
    purge_cache()
    with _lock:
        entry = _cache.get(token)
    if not entry:
        abort(404, "Link expired or not found.")
    data, filename, expiry = entry
    if datetime.now() > expiry:
        with _lock:
            _cache.pop(token, None)
        abort(410, "Link expired (10 min). Please regenerate.")
    return send_file(
        io.BytesIO(data),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename,
    )


@app.route("/health")
def health():
    return jsonify({
        "status":           "ok",
        "template_present": os.path.exists(TEMPLATE_PATH),
        "cache_entries":    len(_cache),
    })


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)
